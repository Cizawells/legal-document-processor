from fastapi import FastAPI, Form, HTTPException, status
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse
from pydantic import BaseModel, Field
import fitz  # PyMuPDF
import json
import logging
import os
from pathlib import Path
from typing import List, Dict, Optional
import re
from uuid import uuid4
from datetime import datetime
from docx import Document
from docx.shared import Pt, RGBColor, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.enum.text import PP_ALIGN
from io import BytesIO
import boto3
from botocore.exceptions import ClientError, NoCredentialsError
from dotenv import load_dotenv
from pii_detector import PIIDetector

# Load environment variables from .env file
load_dotenv()

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

# Storage Configuration
STORAGE_TYPE = os.getenv("STORAGE_TYPE", "local")  # 'local', 's3', or 'r2'
AWS_REGION = os.getenv("AWS_REGION", "us-west-2")
AWS_ACCESS_KEY_ID = os.getenv("AWS_ACCESS_KEY_ID")
AWS_SECRET_ACCESS_KEY = os.getenv("AWS_SECRET_ACCESS_KEY")
S3_BUCKET_NAME = os.getenv("S3_BUCKET_NAME", "pdf-app-files")

# Cloudflare R2 Configuration
R2_ACCESS_KEY_ID = os.getenv("R2_ACCESS_KEY_ID")
R2_SECRET_ACCESS_KEY = os.getenv("R2_SECRET_ACCESS_KEY")
R2_BUCKET_NAME = os.getenv("R2_BUCKET_NAME", "pdf-app-files")
R2_ACCOUNT_ID = os.getenv("R2_ACCOUNT_ID")
R2_ENDPOINT = os.getenv("R2_ENDPOINT")

# Local storage paths (fallback and temp processing)
UPLOADS_PATH = Path(os.getenv("UPLOADS_PATH", "../pdf-merger-backed/uploads")).resolve()
TEMP_PATH = Path(os.getenv("TEMP_PATH", "../pdf-merger-backed/temp")).resolve()

# Initialize cloud storage client (S3 or R2)
s3_client = None
bucket_name = None

if STORAGE_TYPE == "s3":
    if not AWS_ACCESS_KEY_ID or not AWS_SECRET_ACCESS_KEY:
        logger.warning("S3 storage configured but AWS credentials not found. Falling back to local storage.")
        STORAGE_TYPE = "local"
    else:
        try:
            s3_client = boto3.client(
                's3',
                region_name=AWS_REGION,
                aws_access_key_id=AWS_ACCESS_KEY_ID,
                aws_secret_access_key=AWS_SECRET_ACCESS_KEY
            )
            bucket_name = S3_BUCKET_NAME
            logger.info(f"S3 client initialized for bucket: {bucket_name}")
        except Exception as e:
            logger.error(f"Failed to initialize S3 client: {e}")
            STORAGE_TYPE = "local"

elif STORAGE_TYPE == "r2":
    if not R2_ACCESS_KEY_ID or not R2_SECRET_ACCESS_KEY or not R2_ACCOUNT_ID:
        logger.warning("R2 storage configured but credentials not found. Falling back to local storage.")
        STORAGE_TYPE = "local"
    else:
        try:
            # R2 endpoint URL
            endpoint_url = R2_ENDPOINT or f"https://{R2_ACCOUNT_ID}.r2.cloudflarestorage.com"
            
            s3_client = boto3.client(
                's3',
                region_name='auto',  # R2 uses 'auto' as region
                aws_access_key_id=R2_ACCESS_KEY_ID,
                aws_secret_access_key=R2_SECRET_ACCESS_KEY,
                endpoint_url=endpoint_url
            )
            bucket_name = R2_BUCKET_NAME
            logger.info(f"Cloudflare R2 client initialized for bucket: {bucket_name}")
        except Exception as e:
            logger.error(f"Failed to initialize R2 client: {e}")
            STORAGE_TYPE = "local"

# Ensure local directories exist (needed for temp processing even with S3)
UPLOADS_PATH.mkdir(parents=True, exist_ok=True)
TEMP_PATH.mkdir(parents=True, exist_ok=True)

logger.info(f"Storage type: {STORAGE_TYPE}")
if STORAGE_TYPE in ["s3", "r2"]:
    logger.info(f"Cloud storage bucket: {bucket_name}")
else:
    logger.info(f"Local uploads: {UPLOADS_PATH}")
    logger.info(f"Local temp: {TEMP_PATH}")

# Configuration
MAX_FILE_SIZE_MB = 50
MAX_REDACTIONS = 1000

app = FastAPI(
    title="PDF Redaction Microservice",
    version="1.0.0",
    description="Secure PDF redaction service using PyMuPDF"
)

# CORS Configuration - CRITICAL for cross-service communication
allowed_origins = os.getenv("ALLOWED_ORIGINS", "").split(",")
if not allowed_origins or allowed_origins == [""]:
    # Default to localhost for development
    allowed_origins = ["http://localhost:3000", "http://localhost:5001"]

app.add_middleware(
    CORSMiddleware,
    allow_origins=allowed_origins,  # Your NestJS backend URL
    allow_credentials=True,
    allow_methods=["GET", "POST", "PUT", "DELETE", "OPTIONS"],
    allow_headers=["*"],
)

# Pydantic Models
class MergeRequest(BaseModel):
    fileIds: List[str] = Field(..., min_length=2, description="List of PDF file IDs to merge")
    outputName: Optional[str] = Field(None, description="Optional output filename")

class ConvertRequest(BaseModel):
    fileId: str = Field(..., description="PDF file ID to convert")
    outputName: Optional[str] = Field(None, description="Optional output filename")

class CompressRequest(BaseModel):
    fileId: str = Field(..., description="PDF file ID to compress")
    compressionLevel: str = Field(..., description="Compression level: low, medium, high")
    outputName: Optional[str] = Field(None, description="Optional output filename")

class PIIDetectionRequest(BaseModel):
    fileId: str = Field(..., description="PDF file ID to analyze")
    categories: List[str] = Field(..., description="List of PII categories to detect")
    confidenceThreshold: float = Field(0.7, ge=0.0, le=1.0, description="Minimum confidence score")

class RedactionRequest(BaseModel):
    fileId: str = Field(..., description="PDF file ID to redact")
    outputName: Optional[str] = Field("redacted.pdf", description="Output filename")
    areas: List[Dict] = Field(..., description="List of redaction areas")
    settings: Optional[Dict] = Field(None, description="Redaction settings")

class SplitByPatternRequest(BaseModel):
    fileId: str = Field(..., description="PDF file ID to split")
    splitByPattern: str = Field(..., description="Split pattern (e.g., '2' for every 2 pages)")
    outputName: Optional[str] = Field(None, description="Base output filename")

class SplitByRangeRequest(BaseModel):
    fileId: str = Field(..., description="PDF file ID to split")
    splitByRange: str = Field(..., description="Page ranges (e.g., '1-5,6-10,11-15')")
    outputName: Optional[str] = Field(None, description="Base output filename")

class ExtractPagesRequest(BaseModel):
    fileId: str = Field(..., description="PDF file ID to extract from")
    extractPages: str = Field(..., description="Pages to extract (e.g., '1,3,5-7,10')")
    outputName: Optional[str] = Field(None, description="Base output filename")

class SplitBySizeRequest(BaseModel):
    fileId: str = Field(..., description="PDF file ID to split")
    maxSizeKB: int = Field(..., description="Maximum file size in KB")
    outputName: Optional[str] = Field(None, description="Base output filename")

def download_from_s3(file_id: str, folder: str = "uploads") -> Path:
    """Download file from cloud storage (S3/R2) to local temp directory for processing."""
    if not s3_client:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail="Cloud storage client not initialized"
        )
    
    s3_key = f"{folder}/{file_id}"
    local_path = TEMP_PATH / f"temp_{file_id}"
    
    try:
        logger.info(f"Downloading from cloud storage: bucket={bucket_name}, key={s3_key} -> {local_path}")
        print(f"DEBUG: Attempting cloud storage download - bucket={bucket_name}, key={s3_key}")
        s3_client.download_file(bucket_name, s3_key, str(local_path))
        return local_path
    except ClientError as e:
        error_code = e.response['Error']['Code']
        if error_code == 'NoSuchKey':
            raise HTTPException(
                status_code=status.HTTP_404_NOT_FOUND,
                detail=f"File not found in cloud storage: {file_id}"
            )
        else:
            logger.error(f"Cloud storage download error: {e}")
            raise HTTPException(
                status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
                detail=f"Failed to download file from cloud storage: {error_code}"
            )

def upload_to_s3(local_path: Path, file_id: str, folder: str = "temp") -> str:
    """Upload file from local temp directory to cloud storage (S3/R2)."""
    if not s3_client:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail="Cloud storage client not initialized"
        )
    
    s3_key = f"{folder}/{file_id}"
    
    try:
        logger.info(f"Uploading to cloud storage: {local_path} -> {s3_key}")
        s3_client.upload_file(str(local_path), bucket_name, s3_key)
        
        # Clean up local temp file
        local_path.unlink(missing_ok=True)
        
        return s3_key
    except Exception as e:
        logger.error(f"Cloud storage upload error: {e}")
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Failed to upload file to cloud storage: {str(e)}"
        )

def get_file_path(file_id: str, folder: str = "uploads") -> Path:
    """Get file path, downloading from cloud storage (S3/R2) if necessary."""
    if STORAGE_TYPE in ["s3", "r2"]:
        return download_from_s3(file_id, folder)
    else:
        local_path = UPLOADS_PATH / file_id if folder == "uploads" else TEMP_PATH / file_id
        if not local_path.exists():
            raise HTTPException(
                status_code=status.HTTP_404_NOT_FOUND,
                detail=f"File not found: {file_id}"
            )
        return local_path

def save_output_file(local_path: Path, file_id: str) -> str:
    """Save output file to appropriate storage (S3/R2 or local)."""
    if STORAGE_TYPE in ["s3", "r2"]:
        upload_to_s3(local_path, file_id, "temp")
        return file_id
    else:
        # File is already in the correct location for local storage
        return file_id

def sanitize_filename(filename: str) -> str:
    """
    Sanitize filename to prevent path traversal attacks.
    Only allows alphanumeric, dash, underscore, and dot.
    """
    # Remove any path separators
    filename = os.path.basename(filename)
    # Only allow safe characters
    if not re.match(r'^[\w\-\.]+$', filename):
        raise ValueError(f"Invalid filename: {filename}")
    return filename

def validate_redaction_area(area: dict, page_count: int) -> bool:
    """Validate redaction area parameters."""
    try:
        page = area.get("page", 1)
        x = float(area.get("x", 0))
        y = float(area.get("y", 0))
        width = float(area.get("width", 0))
        height = float(area.get("height", 0))
        
        # Validate page number
        if page < 1 or page > page_count:
            return False
        
        # Validate coordinates (allow reasonable bounds)
        if x < 0 or y < 0 or width <= 0 or height <= 0:
            return False
        
        # Prevent extremely large areas (potential DoS)
        if width > 10000 or height > 10000:
            return False
            
        return True
    except (ValueError, TypeError):
        return False

@app.get("/")
async def root():
    """Root endpoint."""
    return {
        "service": "PDF Redaction Microservice",
        "status": "running",
        "version": "1.0.0"
    }

@app.get("/health")
async def health_check():
    """Health check endpoint for monitoring."""
    return {
        "status": "healthy",
        "timestamp": datetime.utcnow().isoformat(),
        "uploads_path": str(UPLOADS_PATH),
        "uploads_path_exists": UPLOADS_PATH.exists(),
        "temp_path": str(TEMP_PATH),
        "temp_path_exists": TEMP_PATH.exists()
    }

@app.post("/redact")
async def redact_pdf(request: RedactionRequest):
    """
    Redacts a PDF stored in shared storage.
    Returns only the output filename (not the file itself).
    """
    try:
        logger.info(f"Redaction request received for file: {request.fileId}")
        
        # Sanitize filenames to prevent path traversal
        try:
            safe_file_id = sanitize_filename(request.fileId)
            safe_output_name = sanitize_filename(request.outputName)
        except ValueError as e:
            logger.warning(f"Invalid filename detected: {e}")
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="Invalid filename format"
            )
        
        # Get file path (downloads from S3 if necessary)
        input_path = get_file_path(safe_file_id, "uploads")
        
        # Check file size
        file_size_mb = input_path.stat().st_size / (1024 * 1024)
        if file_size_mb > MAX_FILE_SIZE_MB:
            logger.warning(f"File too large: {file_size_mb}MB")
            raise HTTPException(
                status_code=status.HTTP_413_REQUEST_ENTITY_TOO_LARGE,
                detail=f"File size exceeds {MAX_FILE_SIZE_MB}MB limit"
            )
        
        # Validate redaction areas
        areas_list = request.areas
        if not isinstance(areas_list, list):
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="Areas must be a list"
            )
        
        # Validate redaction count
        if len(areas_list) > MAX_REDACTIONS:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail=f"Too many redactions (max: {MAX_REDACTIONS})"
            )
        
        # Open PDF
        try:
            doc = fitz.open(str(input_path))
        except Exception as e:
            logger.error(f"Failed to open PDF: {e}")
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="Invalid or corrupted PDF file"
            )
        
        try:
            page_count = len(doc)
            redactions_applied = 0
            
            # Apply redactions
            for item in areas_list:
                logger.info(f"Processing redaction area: {item}")
                
                # Validate redaction area
                if not validate_redaction_area(item, page_count):
                    logger.warning(f"Invalid redaction area: {item}")
                    continue
                
                page_number = item.get("page", 1) - 1  # adjust to 0-based index
                logger.info(f"Applying redaction on page {page_number} (0-based) from original page {item.get('page', 1)} (1-based)")
                
                page = doc[page_number]
                page_rect = page.rect
                logger.info(f"Page {page_number} dimensions: {page_rect.width} x {page_rect.height}")
                
                rect = fitz.Rect(
                    float(item["x"]),
                    float(item["y"]),
                    float(item["x"]) + float(item["width"]),
                    float(item["y"]) + float(item["height"])
                )
                
                logger.info(f"Redaction rectangle: {rect}")
                logger.info(f"Rectangle bounds: x1={rect.x0}, y1={rect.y0}, x2={rect.x1}, y2={rect.y1}")
                
                # Validate rectangle is within page bounds
                if rect.x0 < 0 or rect.y0 < 0 or rect.x1 > page_rect.width or rect.y1 > page_rect.height:
                    logger.warning(f"Rectangle extends beyond page bounds. Page: {page_rect}, Rect: {rect}")
                
                # Add redaction annotation
                page.add_redact_annot(rect, fill=(0, 0, 0))
                redactions_applied += 1
                logger.info(f"Successfully added redaction annotation {redactions_applied}")
            
            # Apply all redactions
            for page_num in range(page_count):
                doc[page_num].apply_redactions()
            
            # Save output file in temp folder
            output_path = TEMP_PATH / safe_output_name
            doc.save(str(output_path))
            
            # Save to appropriate storage (S3 or local)
            final_file_id = save_output_file(output_path, safe_output_name)
            
            logger.info(f"Redacted PDF saved: {safe_output_name} ({redactions_applied} areas)")
            
            return {
                "status": "success",
                "fileId": final_file_id,
                "fileName": final_file_id,
                "redactions_applied": redactions_applied,
                "timestamp": datetime.utcnow().isoformat()
            }
            
        finally:
            # Always close the document
            doc.close()
            
        # Note: File cleanup is handled by the centralized NestJS file tracking service
            
    except HTTPException:
        # Re-raise HTTP exceptions
        raise
    except Exception as e:
        # Catch-all for unexpected errors
        logger.error(f"Unexpected error during redaction: {e}", exc_info=True)
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail="Internal server error during redaction"
        )

@app.post("/merge")
async def merge_pdfs(request: MergeRequest):
    """
    Merge multiple PDFs stored in shared storage using PyMuPDF.

    Expects JSON body:
    {
      "fileIds": ["a.pdf", "b.pdf", ...],
      "outputName": "optional-name.pdf"  // optional
    }

    Returns the output filename (saved under TEMP_PATH) and metadata.
    """
    out_doc = None
    print("merggggin")
    try:
        logger.info(f"Merge request received for {len(request.fileIds)} files")

        # Sanitize and resolve input paths
        source_paths: List[Path] = []
        for fid in request.fileIds:
            try:
                safe_id = sanitize_filename(fid)
            except ValueError as e:
                logger.warning(f"Invalid filename in merge request: {e}")
                raise HTTPException(
                    status_code=status.HTTP_400_BAD_REQUEST,
                    detail=f"Invalid filename: {fid}"
                )
            
            # Get file path (downloads from S3 if necessary)
            input_path = get_file_path(safe_id, "uploads")
            
            # Check file size
            file_size_mb = input_path.stat().st_size / (1024 * 1024)
            if file_size_mb > MAX_FILE_SIZE_MB:
                logger.warning(f"File too large in merge: {file_size_mb}MB")
                raise HTTPException(
                    status_code=status.HTTP_413_REQUEST_ENTITY_TOO_LARGE,
                    detail=f"File {safe_id} exceeds {MAX_FILE_SIZE_MB}MB limit"
                )
            
            source_paths.append(input_path)

        # Determine output name
        if request.outputName:
            try:
                safe_output_name = sanitize_filename(request.outputName)
            except ValueError:
                logger.warning(f"Invalid output name: {request.outputName}")
                raise HTTPException(
                    status_code=status.HTTP_400_BAD_REQUEST,
                    detail="Invalid outputName format"
                )
        else:
            safe_output_name = f"merged-{uuid4().hex[:8]}.pdf"

        output_path = TEMP_PATH / safe_output_name

        # Perform merge
        merged_pages = 0
        out_doc = fitz.open()  # Create empty PDF
        
        for src_path in source_paths:
            try:
                src_doc = fitz.open(str(src_path))
                out_doc.insert_pdf(src_doc)  # Insert all pages
                merged_pages += len(src_doc)
                src_doc.close()
            except Exception as e:
                logger.error(f"Failed to process PDF {src_path.name}: {e}")
                raise HTTPException(
                    status_code=status.HTTP_400_BAD_REQUEST,
                    detail=f"Invalid or corrupted PDF: {src_path.name}"
                )
        
        # Save merged document
        out_doc.save(str(output_path))
        
        # Save to appropriate storage (S3 or local)
        final_file_id = save_output_file(output_path, safe_output_name)
        
        logger.info(
            f"Successfully merged {len(source_paths)} PDFs into {safe_output_name} "
            f"with {merged_pages} pages"
        )

        return {
            "status": "success",
            "fileId": final_file_id,
            "sourceCount": len(source_paths),
            "pageCount": merged_pages,
            "timestamp": datetime.utcnow().isoformat(),
        }

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Unexpected error during merge: {e}", exc_info=True)
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail="Internal server error during merge"
        )
    finally:
        # Always close the output document if it was created
        if out_doc is not None:
            try:
                out_doc.close()
            except Exception as e:
                logger.warning(f"Failed to close output document: {e}")
        
        # Note: File cleanup is handled by the centralized NestJS file tracking service

@app.post("/convert/pdf-to-word")
async def convert_pdf_to_word(request: ConvertRequest):
    """
    Convert PDF to Word document using PyMuPDF text extraction.
    
    Args:
        request: ConvertRequest with fileId and optional outputName
        
    Returns:
        JSON with status, fileId (output filename), pageCount, and timestamp
    """

    print("converttting to word", request)
    pdf_doc = None
    try:
        logger.info(f"PDF to Word conversion request received for file: {request.fileId}")
        
        # Sanitize and validate input file
        try:
            safe_id = sanitize_filename(request.fileId)
        except ValueError as e:
            logger.warning(f"Invalid filename in convert request: {e}")
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail=f"Invalid filename: {request.fileId}"
            )
        
        # Get file path (downloads from S3 if necessary)
        print(f"DEBUG: About to get file path for fileId: {safe_id}")
        print(f"DEBUG: STORAGE_TYPE: {STORAGE_TYPE}")
        print(f"DEBUG: BUCKET_NAME: {bucket_name}")
        input_path = get_file_path(safe_id, "uploads")
        
        # Check file size
        file_size_mb = input_path.stat().st_size / (1024 * 1024)
        if file_size_mb > MAX_FILE_SIZE_MB:
            logger.warning(f"File too large for conversion: {file_size_mb}MB")
            raise HTTPException(
                status_code=status.HTTP_413_REQUEST_ENTITY_TOO_LARGE,
                detail=f"File {safe_id} exceeds {MAX_FILE_SIZE_MB}MB limit"
            )
        
        # Determine output name
        if request.outputName:
            try:
                safe_output_name = sanitize_filename(request.outputName)
                # Ensure .docx extension
                if not safe_output_name.endswith('.docx'):
                    safe_output_name = safe_output_name.rsplit('.', 1)[0] + '.docx'
            except ValueError:
                logger.warning(f"Invalid output name: {request.outputName}")
                raise HTTPException(
                    status_code=status.HTTP_400_BAD_REQUEST,
                    detail="Invalid outputName format"
                )
        else:
            safe_output_name = f"converted-{uuid4().hex[:8]}.docx"
        
        output_path = TEMP_PATH / safe_output_name
        
        # Open PDF and extract content
        pdf_doc = fitz.open(str(input_path))
        page_count = len(pdf_doc)
        
        # Create Word document
        doc = Document()
        
        # Set document margins
        sections = doc.sections
        for section in sections:
            section.top_margin = Inches(1)
            section.bottom_margin = Inches(1)
            section.left_margin = Inches(1)
            section.right_margin = Inches(1)
        
        # Extract text from each page
        for page_num in range(page_count):
            page = pdf_doc[page_num]
            
            # Extract text with layout preservation
            text = page.get_text("text")
            
            # Add page number heading
            if page_num > 0:
                doc.add_page_break()
            
            heading = doc.add_paragraph(f"Page {page_num + 1}")
            heading.style = 'Heading 2'
            
            # Add text content
            if text.strip():
                # Split into paragraphs and add to document
                paragraphs = text.split('\n\n')
                for para_text in paragraphs:
                    if para_text.strip():
                        para = doc.add_paragraph(para_text.strip())
                        para.style = 'Normal'
            else:
                # Page has no text
                para = doc.add_paragraph("[No text content on this page]")
                para.style = 'Normal'
                run = para.runs[0]
                run.italic = True
                run.font.color.rgb = RGBColor(128, 128, 128)
        
        # Save Word document
        doc.save(str(output_path))
        
        # Save to appropriate storage (S3 or local)
        final_file_id = save_output_file(output_path, safe_output_name)
        
        logger.info(
            f"Successfully converted PDF to Word: {safe_output_name} "
            f"with {page_count} pages"
        )
        
        return {
            "status": "success",
            "fileName": final_file_id,
            "pageCount": page_count,
            "timestamp": datetime.utcnow().isoformat(),
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Unexpected error during PDF to Word conversion: {e}", exc_info=True)
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail="Internal server error during conversion"
        )
    finally:
        # Always close the PDF document if it was opened
        if pdf_doc is not None:
            try:
                pdf_doc.close()
            except Exception as e:
                logger.warning(f"Failed to close PDF document: {e}")
        
        # Note: File cleanup is handled by the centralized NestJS file tracking service
@app.post("/convert/pdf-to-powerpoint")
async def convert_pdf_to_powerpoint(request: ConvertRequest):
    """
    Convert PDF to PowerPoint presentation using PyMuPDF text extraction.
    
    Args:
        request: ConvertRequest with fileId and optional outputName
        
    Returns:
        JSON with status, fileName (output filename), pageCount, and timestamp
    """
    pdf_doc = None
    try:
        logger.info(f"PDF to PowerPoint conversion request received for file: {request.fileId}")
        
        # Sanitize and validate input file
        try:
            safe_id = sanitize_filename(request.fileId)
        except ValueError as e:
            logger.warning(f"Invalid filename in convert request: {e}")
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail=f"Invalid filename: {request.fileId}"
            )
        
        # Get file path (downloads from S3 if necessary)
        input_path = get_file_path(safe_id, "uploads")
        
        # Check file size
        file_size_mb = input_path.stat().st_size / (1024 * 1024)
        if file_size_mb > MAX_FILE_SIZE_MB:
            logger.warning(f"File too large for conversion: {file_size_mb}MB")
            raise HTTPException(
                status_code=status.HTTP_413_REQUEST_ENTITY_TOO_LARGE,
                detail=f"File {safe_id} exceeds {MAX_FILE_SIZE_MB}MB limit"
            )
        
        # Determine output name
        if request.outputName:
            try:
                safe_output_name = sanitize_filename(request.outputName)
                # Ensure .pptx extension
                if not safe_output_name.endswith('.pptx'):
                    safe_output_name = safe_output_name.rsplit('.', 1)[0] + '.pptx'
            except ValueError:
                logger.warning(f"Invalid output name: {request.outputName}")
                raise HTTPException(
                    status_code=status.HTTP_400_BAD_REQUEST,
                    detail="Invalid outputName format"
                )
        else:
            safe_output_name = f"converted-{uuid4().hex[:8]}.pptx"
        
        output_path = TEMP_PATH / safe_output_name
        
        # Open PDF and extract content
        pdf_doc = fitz.open(str(input_path))
        page_count = len(pdf_doc)
        
        # Create PowerPoint presentation
        prs = Presentation()
        
        # Set slide dimensions (16:9 aspect ratio)
        prs.slide_width = PptxInches(10)
        prs.slide_height = PptxInches(7.5)
        
        # Extract text and images from each page
        for page_num in range(page_count):
            page = pdf_doc[page_num]
            
            # Add a blank slide
            blank_slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Extract text with layout preservation
            text = page.get_text("text")
            
            # Add title text box at the top
            title_left = PptxInches(0.5)
            title_top = PptxInches(0.5)
            title_width = PptxInches(9)
            title_height = PptxInches(0.8)
            
            title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
            title_frame = title_box.text_frame
            title_frame.text = f"Page {page_num + 1}"
            title_para = title_frame.paragraphs[0]
            title_para.font.size = PptxPt(24)
            title_para.font.bold = True
            
            # Add content text box
            content_left = PptxInches(0.5)
            content_top = PptxInches(1.5)
            content_width = PptxInches(9)
            content_height = PptxInches(5.5)
            
            content_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            
            if text.strip():
                # Split into paragraphs and add to slide
                paragraphs = text.split('\n\n')
                for i, para_text in enumerate(paragraphs):
                    if para_text.strip():
                        if i == 0:
                            content_frame.text = para_text.strip()
                        else:
                            p = content_frame.add_paragraph()
                            p.text = para_text.strip()
                            p.level = 0
                        
                # Set font size for all paragraphs
                for paragraph in content_frame.paragraphs:
                    paragraph.font.size = PptxPt(14)
            else:
                # Page has no text
                content_frame.text = "[No text content on this page]"
                para = content_frame.paragraphs[0]
                para.font.size = PptxPt(14)
                para.font.italic = True
                para.font.color.rgb = (128, 128, 128)
            
            # Try to extract and add images
            try:
                image_list = page.get_images()
                if image_list:
                    # Add first image if available (to keep slides simple)
                    img_index = image_list[0][0]
                    base_image = pdf_doc.extract_image(img_index)
                    image_bytes = base_image["image"]
                    
                    # Save image temporarily
                    img_stream = BytesIO(image_bytes)
                    
                    # Add image to slide (bottom right corner)
                    img_left = PptxInches(7)
                    img_top = PptxInches(5)
                    img_width = PptxInches(2.5)
                    
                    try:
                        slide.shapes.add_picture(img_stream, img_left, img_top, width=img_width)
                    except Exception as img_err:
                        logger.warning(f"Could not add image to slide {page_num + 1}: {img_err}")
            except Exception as e:
                logger.warning(f"Could not extract images from page {page_num + 1}: {e}")
        
        # Save PowerPoint presentation
        prs.save(str(output_path))
        
        # Save to appropriate storage (S3 or local)
        final_file_id = save_output_file(output_path, safe_output_name)
        
        logger.info(
            f"Successfully converted PDF to PowerPoint: {safe_output_name} "
            f"with {page_count} slides"
        )
        
        return {
            "status": "success",
            "fileName": final_file_id,
            "pageCount": page_count,
            "timestamp": datetime.utcnow().isoformat(),
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Unexpected error during PDF to PowerPoint conversion: {e}", exc_info=True)
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail="Internal server error during conversion"
        )
    finally:
        # Always close the PDF document if it was opened
        if pdf_doc is not None:
            try:
                pdf_doc.close()
            except Exception as e:
                logger.warning(f"Failed to close PDF document: {e}")
        
        # Note: File cleanup is handled by the centralized NestJS file tracking service

@app.post("/compress")
async def compress_pdf(request: CompressRequest):
    """
    Compress PDF using PyMuPDF with various optimization strategies.
    
    Args:
        request: CompressRequest with fileId, compressionLevel, and optional outputName
        
    Returns:
        JSON with status, fileName, originalSize, compressedSize, compressionRatio, and timestamp
    """
    pdf_doc = None
    try:
        logger.info(f"PDF compression request received for file: {request.fileId}, level: {request.compressionLevel}")
        
        # Validate compression level
        valid_levels = ["low", "medium", "high"]
        if request.compressionLevel not in valid_levels:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail=f"Invalid compression level. Must be one of: {valid_levels}"
            )
        
        # Sanitize and validate input file
        try:
            safe_id = sanitize_filename(request.fileId)
        except ValueError as e:
            logger.warning(f"Invalid filename in compress request: {e}")
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail=f"Invalid filename: {request.fileId}"
            )
        
        # Get file path (downloads from S3 if necessary)
        input_path = get_file_path(safe_id, "uploads")
        
        # Check file size
        original_size = input_path.stat().st_size
        file_size_mb = original_size / (1024 * 1024)
        if file_size_mb > MAX_FILE_SIZE_MB:
            logger.warning(f"File too large for compression: {file_size_mb}MB")
            raise HTTPException(
                status_code=status.HTTP_413_REQUEST_ENTITY_TOO_LARGE,
                detail=f"File {safe_id} exceeds {MAX_FILE_SIZE_MB}MB limit"
            )
        
        # Determine output name
        if request.outputName:
            try:
                safe_output_name = sanitize_filename(request.outputName)
                # Ensure .pdf extension
                if not safe_output_name.endswith('.pdf'):
                    safe_output_name = safe_output_name.rsplit('.', 1)[0] + '.pdf'
            except ValueError:
                logger.warning(f"Invalid output name: {request.outputName}")
                raise HTTPException(
                    status_code=status.HTTP_400_BAD_REQUEST,
                    detail="Invalid outputName format"
                )
        else:
            safe_output_name = f"compressed-{uuid4().hex[:8]}.pdf"
        
        output_path = TEMP_PATH / safe_output_name
        
        # Open PDF and perform compression
        pdf_doc = fitz.open(str(input_path))
        
        # Set compression parameters based on level
        compression_params = {
            "high": {"deflate_images": True, "deflate_fonts": True, "garbage": 3, "image_quality": 85},
            "medium": {"deflate_images": True, "deflate_fonts": True, "garbage": 4, "image_quality": 75},
            "low": {"deflate_images": True, "deflate_fonts": True, "garbage": 4, "image_quality": 60}
        }
        
        params = compression_params[request.compressionLevel]
        
        # Compress images in the PDF
        for page_num in range(len(pdf_doc)):
            page = pdf_doc[page_num]
            
            # Get all images on the page
            image_list = page.get_images()
            
            for img_index, img in enumerate(image_list):
                try:
                    # Get image data
                    xref = img[0]
                    base_image = pdf_doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    image_ext = base_image["ext"]
                    
                    # Only compress if it's a significant image (> 10KB)
                    if len(image_bytes) > 10240:
                        # Convert to PIL Image for compression
                        from PIL import Image
                        import io
                        
                        # Open image with PIL
                        pil_image = Image.open(io.BytesIO(image_bytes))
                        
                        # Convert to RGB if necessary
                        if pil_image.mode in ('RGBA', 'LA', 'P'):
                            pil_image = pil_image.convert('RGB')
                        
                        # Compress the image
                        compressed_image_bytes = io.BytesIO()
                        pil_image.save(
                            compressed_image_bytes, 
                            format='JPEG', 
                            quality=params["image_quality"],
                            optimize=True
                        )
                        compressed_image_bytes.seek(0)
                        
                        # Replace the image in the PDF
                        pdf_doc.replace_image(xref, stream=compressed_image_bytes.getvalue())
                        
                except Exception as img_err:
                    logger.warning(f"Could not compress image {img_index} on page {page_num + 1}: {img_err}")
                    continue
        
        # Save with comprehensive compression options
        pdf_doc.save(
            str(output_path),
            garbage=params["garbage"],  # Remove unused objects (3=conservative, 4=aggressive)
            deflate=True,  # Compress content streams
            deflate_images=params["deflate_images"],  # Compress images
            deflate_fonts=params["deflate_fonts"],  # Compress fonts
            clean=True,  # Clean up document structure
            ascii=False,  # Use binary encoding
            expand=0,  # Don't expand images
            linear=False,  # Don't linearize
            pretty=False,  # Don't pretty-print
            encryption=fitz.PDF_ENCRYPT_NONE,
            permissions=-1,
            owner_pw=None,
            user_pw=None
        )
        
        # Calculate compression statistics
        compressed_size = output_path.stat().st_size
        compression_ratio = ((original_size - compressed_size) / original_size) * 100
        
        logger.info(
            f"Successfully compressed PDF: {safe_output_name} "
            f"({original_size} -> {compressed_size} bytes, {compression_ratio:.1f}% reduction)"
        )
        
        # Save to appropriate storage (S3 or local)
        final_file_id = save_output_file(output_path, safe_output_name)
        
        return {
            "status": "success",
            "fileName": final_file_id,
            "originalSize": original_size,
            "compressedSize": compressed_size,
            "compressionRatio": round(compression_ratio, 1),
            "compressionLevel": request.compressionLevel,
            "timestamp": datetime.utcnow().isoformat(),
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Unexpected error during PDF compression: {e}", exc_info=True)
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail="Internal server error during compression"
        )
    finally:
        # Always close the PDF document if it was opened
        if pdf_doc is not None:
            try:
                pdf_doc.close()
            except Exception as e:
                logger.warning(f"Failed to close PDF document: {e}")

@app.post("/detect-pii")
async def detect_pii(request: PIIDetectionRequest):
    """
    Detect PII (Personally Identifiable Information) in a PDF document.
    
    Args:
        request: PIIDetectionRequest with fileId, categories, and confidenceThreshold
        
    Returns:
        JSON with findings array and statistics
    """
    pdf_doc = None
    try:
        logger.info(f"PII detection request received for file: {request.fileId}")
        
        # Sanitize filename
        try:
            safe_file_id = sanitize_filename(request.fileId)
        except ValueError as e:
            logger.warning(f"Invalid filename detected: {e}")
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="Invalid filename format"
            )
        
        # Get file path
        input_path = get_file_path(safe_file_id, "uploads")
        
        # Check file size
        file_size_mb = input_path.stat().st_size / (1024 * 1024)
        if file_size_mb > MAX_FILE_SIZE_MB:
            logger.warning(f"File too large: {file_size_mb}MB")
            raise HTTPException(
                status_code=status.HTTP_413_REQUEST_ENTITY_TOO_LARGE,
                detail=f"File size exceeds {MAX_FILE_SIZE_MB}MB limit"
            )
        
        # Open PDF
        try:
            pdf_doc = fitz.open(str(input_path))
        except Exception as e:
            logger.error(f"Failed to open PDF: {e}")
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="Invalid or corrupted PDF file"
            )
        
        findings = []
        page_count = len(pdf_doc)
        
        # Process each page
        for page_num in range(page_count):
            page = pdf_doc[page_num]
            
            # Extract text with position information
            text_instances = page.get_text("dict")
            
            # Process each block
            for block in text_instances.get("blocks", []):
                if block.get("type") == 0:  # Text block
                    for line in block.get("lines", []):
                        for span in line.get("spans", []):
                            text = span.get("text", "")
                            bbox = span.get("bbox", [0, 0, 0, 0])
                            
                            # Detect PII in this text span
                            span_findings = PIIDetector.detect_in_text(
                                text=text,
                                bbox=list(bbox),
                                page_num=page_num + 1,  # 1-based page numbering
                                categories=request.categories,
                                confidence_threshold=request.confidenceThreshold
                            )
                            
                            findings.extend(span_findings)
        
        # Calculate statistics
        statistics = PIIDetector.calculate_statistics(findings)
        
        logger.info(f"PII detection completed: {len(findings)} findings in {page_count} pages")
        
        return {
            "status": "success",
            "findings": findings,
            "statistics": statistics,
            "timestamp": datetime.utcnow().isoformat()
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Unexpected error during PII detection: {e}", exc_info=True)
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail="Internal server error during PII detection"
        )
    finally:
        # Always close the PDF document if it was opened
        if pdf_doc is not None:
            try:
                pdf_doc.close()
            except Exception as e:
                logger.warning(f"Failed to close PDF document: {e}")

                
@app.post("/split/pattern")
async def split_pdf_by_pattern(request: SplitByPatternRequest):
    """
    Split PDF by pattern (every N pages).
    
    Args:
        request: SplitByPatternRequest with fileId, splitByPattern, and optional outputName
        
    Returns:
        JSON with status and list of output files
    """
    pdf_doc = None
    try:
        logger.info(f"PDF split by pattern request received for file: {request.fileId}, pattern: {request.splitByPattern}")
        
        # Sanitize and validate input file
        try:
            safe_id = sanitize_filename(request.fileId)
        except ValueError as e:
            logger.warning(f"Invalid filename in split request: {e}")
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail=f"Invalid filename: {request.fileId}"
            )
        
        # Get file path (downloads from S3 if necessary)
        input_path = get_file_path(safe_id, "uploads")
        
        # Check file size
        file_size_mb = input_path.stat().st_size / (1024 * 1024)
        if file_size_mb > MAX_FILE_SIZE_MB:
            logger.warning(f"File too large for split: {file_size_mb}MB")
            raise HTTPException(
                status_code=status.HTTP_413_REQUEST_ENTITY_TOO_LARGE,
                detail=f"File {safe_id} exceeds {MAX_FILE_SIZE_MB}MB limit"
            )
        
        # Parse split pattern (number of pages per split)
        try:
            pages_per_split = int(request.splitByPattern)
            if pages_per_split <= 0:
                raise ValueError("Pages per split must be positive")
        except ValueError:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="Invalid split pattern. Must be a positive integer."
            )
        
        # Determine base output name
        if request.outputName:
            try:
                safe_output_name = sanitize_filename(request.outputName)
                # Remove extension for base name
                if safe_output_name.endswith('.pdf'):
                    safe_output_name = safe_output_name[:-4]
            except ValueError:
                logger.warning(f"Invalid output name: {request.outputName}")
                raise HTTPException(
                    status_code=status.HTTP_400_BAD_REQUEST,
                    detail="Invalid outputName format"
                )
        else:
            safe_output_name = f"split-pattern-{uuid4().hex[:8]}"
        
        # Open PDF and split
        pdf_doc = fitz.open(str(input_path))
        page_count = len(pdf_doc)
        
        if pages_per_split >= page_count:
            # If pattern is larger than page count, return original file
            output_filename = f"{safe_output_name}.pdf"
            output_path = TEMP_PATH / output_filename
            pdf_doc.save(str(output_path))
            
            # Save to appropriate storage
            final_file_id = save_output_file(output_path, output_filename)
            
            return {
                "status": "success",
                "files": [final_file_id]
            }
        
        output_files = []
        
        # Split into chunks
        for start_page in range(0, page_count, pages_per_split):
            end_page = min(start_page + pages_per_split - 1, page_count - 1)
            
            # Create new document for this chunk
            chunk_doc = fitz.open()
            chunk_doc.insert_pdf(pdf_doc, from_page=start_page, to_page=end_page)
            
            # Generate output filename
            chunk_number = (start_page // pages_per_split) + 1
            output_filename = f"{safe_output_name}-{chunk_number}.pdf"
            output_path = TEMP_PATH / output_filename
            
            # Save chunk
            chunk_doc.save(str(output_path))
            chunk_doc.close()
            
            # Save to appropriate storage
            final_file_id = save_output_file(output_path, output_filename)
            output_files.append(final_file_id)
        
        logger.info(f"Successfully split PDF into {len(output_files)} files")
        
        return {
            "status": "success",
            "files": output_files
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Unexpected error during PDF split by pattern: {e}", exc_info=True)
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail="Internal server error during split"
        )
    finally:
        if pdf_doc is not None:
            try:
                pdf_doc.close()
            except Exception as e:
                logger.warning(f"Failed to close PDF document: {e}")

@app.post("/split/range")
async def split_pdf_by_range(request: SplitByRangeRequest):
    """
    Split PDF by custom page ranges.
    
    Args:
        request: SplitByRangeRequest with fileId, splitByRange, and optional outputName
        
    Returns:
        JSON with status and list of output files
    """
    pdf_doc = None
    try:
        logger.info(f"PDF split by range request received for file: {request.fileId}, ranges: {request.splitByRange}")
        
        # Sanitize and validate input file
        try:
            safe_id = sanitize_filename(request.fileId)
        except ValueError as e:
            logger.warning(f"Invalid filename in split request: {e}")
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail=f"Invalid filename: {request.fileId}"
            )
        
        # Get file path (downloads from S3 if necessary)
        input_path = get_file_path(safe_id, "uploads")
        
        # Check file size
        file_size_mb = input_path.stat().st_size / (1024 * 1024)
        if file_size_mb > MAX_FILE_SIZE_MB:
            logger.warning(f"File too large for split: {file_size_mb}MB")
            raise HTTPException(
                status_code=status.HTTP_413_REQUEST_ENTITY_TOO_LARGE,
                detail=f"File {safe_id} exceeds {MAX_FILE_SIZE_MB}MB limit"
            )
        
        # Parse page ranges (e.g., "1-5,6-10,11-15")
        try:
            ranges = []
            for range_str in request.splitByRange.split(','):
                range_str = range_str.strip()
                if '-' in range_str:
                    start, end = range_str.split('-', 1)
                    start_page = int(start.strip()) - 1  # Convert to 0-based
                    end_page = int(end.strip()) - 1      # Convert to 0-based
                    if start_page < 0 or end_page < start_page:
                        raise ValueError(f"Invalid range: {range_str}")
                    ranges.append((start_page, end_page))
                else:
                    # Single page
                    page_num = int(range_str) - 1  # Convert to 0-based
                    if page_num < 0:
                        raise ValueError(f"Invalid page: {range_str}")
                    ranges.append((page_num, page_num))
        except ValueError as e:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail=f"Invalid range format: {str(e)}"
            )
        
        # Determine base output name
        if request.outputName:
            try:
                safe_output_name = sanitize_filename(request.outputName)
                # Remove extension for base name
                if safe_output_name.endswith('.pdf'):
                    safe_output_name = safe_output_name[:-4]
            except ValueError:
                logger.warning(f"Invalid output name: {request.outputName}")
                raise HTTPException(
                    status_code=status.HTTP_400_BAD_REQUEST,
                    detail="Invalid outputName format"
                )
        else:
            safe_output_name = f"split-range-{uuid4().hex[:8]}"
        
        # Open PDF and split
        pdf_doc = fitz.open(str(input_path))
        page_count = len(pdf_doc)
        
        output_files = []
        
        # Process each range
        for i, (start_page, end_page) in enumerate(ranges):
            # Validate range against document
            if start_page >= page_count or end_page >= page_count:
                logger.warning(f"Range {start_page+1}-{end_page+1} exceeds document pages ({page_count})")
                continue
            
            # Create new document for this range
            range_doc = fitz.open()
            range_doc.insert_pdf(pdf_doc, from_page=start_page, to_page=end_page)
            
            # Generate output filename
            if start_page == end_page:
                output_filename = f"{safe_output_name}-page{start_page+1}.pdf"
            else:
                output_filename = f"{safe_output_name}-pages{start_page+1}-{end_page+1}.pdf"
            output_path = TEMP_PATH / output_filename
            
            # Save range
            range_doc.save(str(output_path))
            range_doc.close()
            
            # Save to appropriate storage
            final_file_id = save_output_file(output_path, output_filename)
            output_files.append(final_file_id)
        
        logger.info(f"Successfully split PDF into {len(output_files)} files by ranges")
        
        return {
            "status": "success",
            "files": output_files
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Unexpected error during PDF split by range: {e}", exc_info=True)
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail="Internal server error during split"
        )
    finally:
        if pdf_doc is not None:
            try:
                pdf_doc.close()
            except Exception as e:
                logger.warning(f"Failed to close PDF document: {e}")

@app.post("/split/extract")
async def extract_pdf_pages(request: ExtractPagesRequest):
    """
    Extract specific pages from PDF.
    
    Args:
        request: ExtractPagesRequest with fileId, extractPages, and optional outputName
        
    Returns:
        JSON with status and list of output files
    """
    pdf_doc = None
    try:
        logger.info(f"PDF extract pages request received for file: {request.fileId}, pages: {request.extractPages}")
        
        # Sanitize and validate input file
        try:
            safe_id = sanitize_filename(request.fileId)
        except ValueError as e:
            logger.warning(f"Invalid filename in extract request: {e}")
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail=f"Invalid filename: {request.fileId}"
            )
        
        # Get file path (downloads from S3 if necessary)
        input_path = get_file_path(safe_id, "uploads")
        
        # Check file size
        file_size_mb = input_path.stat().st_size / (1024 * 1024)
        if file_size_mb > MAX_FILE_SIZE_MB:
            logger.warning(f"File too large for extract: {file_size_mb}MB")
            raise HTTPException(
                status_code=status.HTTP_413_REQUEST_ENTITY_TOO_LARGE,
                detail=f"File {safe_id} exceeds {MAX_FILE_SIZE_MB}MB limit"
            )
        
        # Parse pages to extract (e.g., "1,3,5-7,10")
        try:
            pages_to_extract = []
            for page_str in request.extractPages.split(','):
                page_str = page_str.strip()
                if '-' in page_str:
                    # Range of pages
                    start, end = page_str.split('-', 1)
                    start_page = int(start.strip()) - 1  # Convert to 0-based
                    end_page = int(end.strip()) - 1      # Convert to 0-based
                    if start_page < 0 or end_page < start_page:
                        raise ValueError(f"Invalid range: {page_str}")
                    pages_to_extract.extend(range(start_page, end_page + 1))
                else:
                    # Single page
                    page_num = int(page_str) - 1  # Convert to 0-based
                    if page_num < 0:
                        raise ValueError(f"Invalid page: {page_str}")
                    pages_to_extract.append(page_num)
            
            # Remove duplicates and sort
            pages_to_extract = sorted(list(set(pages_to_extract)))
            
        except ValueError as e:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail=f"Invalid page format: {str(e)}"
            )
        
        # Determine output name
        if request.outputName:
            try:
                safe_output_name = sanitize_filename(request.outputName)
                # Ensure .pdf extension
                if not safe_output_name.endswith('.pdf'):
                    safe_output_name += '.pdf'
            except ValueError:
                logger.warning(f"Invalid output name: {request.outputName}")
                raise HTTPException(
                    status_code=status.HTTP_400_BAD_REQUEST,
                    detail="Invalid outputName format"
                )
        else:
            safe_output_name = f"extracted-{uuid4().hex[:8]}.pdf"
        
        # Open PDF and extract pages
        pdf_doc = fitz.open(str(input_path))
        page_count = len(pdf_doc)
        
        # Validate page numbers
        valid_pages = [p for p in pages_to_extract if p < page_count]
        if not valid_pages:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="No valid pages to extract"
            )
        
        # Create new document with extracted pages
        extracted_doc = fitz.open()
        for page_num in valid_pages:
            extracted_doc.insert_pdf(pdf_doc, from_page=page_num, to_page=page_num)
        
        # Save extracted document
        output_path = TEMP_PATH / safe_output_name
        extracted_doc.save(str(output_path))
        extracted_doc.close()
        
        # Save to appropriate storage
        final_file_id = save_output_file(output_path, safe_output_name)
        
        logger.info(f"Successfully extracted {len(valid_pages)} pages from PDF")
        
        return {
            "status": "success",
            "files": [final_file_id]
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Unexpected error during PDF page extraction: {e}", exc_info=True)
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail="Internal server error during extraction"
        )
    finally:
        if pdf_doc is not None:
            try:
                pdf_doc.close()
            except Exception as e:
                logger.warning(f"Failed to close PDF document: {e}")

@app.post("/split/size")
async def split_pdf_by_size(request: SplitBySizeRequest):
    """
    Split PDF by maximum file size.
    
    Args:
        request: SplitBySizeRequest with fileId, maxSizeKB, and optional outputName
        
    Returns:
        JSON with status and list of output files
    """
    pdf_doc = None
    try:
        logger.info(f"PDF split by size request received for file: {request.fileId}, max size: {request.maxSizeKB}KB")
        
        # Sanitize and validate input file
        try:
            safe_id = sanitize_filename(request.fileId)
        except ValueError as e:
            logger.warning(f"Invalid filename in split request: {e}")
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail=f"Invalid filename: {request.fileId}"
            )
        
        # Get file path (downloads from S3 if necessary)
        input_path = get_file_path(safe_id, "uploads")
        
        # Check file size
        file_size_mb = input_path.stat().st_size / (1024 * 1024)
        if file_size_mb > MAX_FILE_SIZE_MB:
            logger.warning(f"File too large for split: {file_size_mb}MB")
            raise HTTPException(
                status_code=status.HTTP_413_REQUEST_ENTITY_TOO_LARGE,
                detail=f"File {safe_id} exceeds {MAX_FILE_SIZE_MB}MB limit"
            )
        
        # Validate max size
        if request.maxSizeKB <= 0:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="Maximum size must be positive"
            )
        
        max_size_bytes = request.maxSizeKB * 1024
        
        # Determine base output name
        if request.outputName:
            try:
                safe_output_name = sanitize_filename(request.outputName)
                # Remove extension for base name
                if safe_output_name.endswith('.pdf'):
                    safe_output_name = safe_output_name[:-4]
            except ValueError:
                logger.warning(f"Invalid output name: {request.outputName}")
                raise HTTPException(
                    status_code=status.HTTP_400_BAD_REQUEST,
                    detail="Invalid outputName format"
                )
        else:
            safe_output_name = f"split-size-{uuid4().hex[:8]}"
        
        # Open PDF and split by size
        pdf_doc = fitz.open(str(input_path))
        page_count = len(pdf_doc)
        
        output_files = []
        current_doc = fitz.open()
        current_pages = 0
        chunk_number = 1
        
        for page_num in range(page_count):
            # Add page to current document
            current_doc.insert_pdf(pdf_doc, from_page=page_num, to_page=page_num)
            current_pages += 1
            
            # Save current document to check size
            temp_output_path = TEMP_PATH / f"temp_size_check_{uuid4().hex[:8]}.pdf"
            current_doc.save(str(temp_output_path))
            current_size = temp_output_path.stat().st_size
            
            # If size exceeds limit or this is the last page, save the chunk
            if current_size > max_size_bytes or page_num == page_count - 1:
                if current_size > max_size_bytes and current_pages > 1:
                    # Remove the last page that caused the overflow
                    current_doc.close()
                    current_doc = fitz.open()
                    for p in range(page_num - current_pages + 1, page_num):
                        current_doc.insert_pdf(pdf_doc, from_page=p, to_page=p)
                    current_doc.save(str(temp_output_path))
                    
                    # Save this chunk
                    output_filename = f"{safe_output_name}-{chunk_number}.pdf"
                    final_output_path = TEMP_PATH / output_filename
                    current_doc.save(str(final_output_path))
                    current_doc.close()
                    
                    # Save to appropriate storage
                    final_file_id = save_output_file(final_output_path, output_filename)
                    output_files.append(final_file_id)
                    
                    # Start new chunk with the overflow page
                    current_doc = fitz.open()
                    current_doc.insert_pdf(pdf_doc, from_page=page_num, to_page=page_num)
                    current_pages = 1
                    chunk_number += 1
                else:
                    # Save current chunk
                    output_filename = f"{safe_output_name}-{chunk_number}.pdf"
                    final_output_path = TEMP_PATH / output_filename
                    current_doc.save(str(final_output_path))
                    current_doc.close()
                    
                    # Save to appropriate storage
                    final_file_id = save_output_file(final_output_path, output_filename)
                    output_files.append(final_file_id)
                    
                    # Reset for next chunk
                    current_doc = fitz.open()
                    current_pages = 0
                    chunk_number += 1
            
            # Clean up temp file
            temp_output_path.unlink(missing_ok=True)
        
        # Close any remaining document
        if current_doc is not None:
            current_doc.close()
        
        logger.info(f"Successfully split PDF into {len(output_files)} files by size")
        
        return {
            "status": "success",
            "files": output_files
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Unexpected error during PDF split by size: {e}", exc_info=True)
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail="Internal server error during split"
        )
    finally:
        if pdf_doc is not None:
            try:
                pdf_doc.close()
            except Exception as e:
                logger.warning(f"Failed to close PDF document: {e}")
        
        # Note: File cleanup is handled by the centralized NestJS file tracking service